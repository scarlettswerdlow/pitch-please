"""Module with util functions for working with pptx Python library"""

from pptx import Presentation
from pptx.util import Inches
import json

SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_CONTENT = 1
SLIDE_LAYOUT_SECTION_HEADER = 2
SLIDE_LAYOUT_TWO_CONTENT = 3
SLIDE_LAYOUT_COMPARISON = 4
SLIDE_LAYOUT_TITLE_ONLY = 5
SLIDE_LAYOUT_BLANK = 6
SLIDE_LAYOUT_CONTENT_CAPTION = 7
SLIDE_LAYOUT_PICTURE_CAPTION = 8

def add_logo_slide(prs, logo_fp:str):
    """Function adding logo slide to presentation"""
    logo_width = logo_height = Inches(6)
    logo_left = (prs.slide_width - logo_width) / 2
    logo_top = (prs.slide_height - logo_height) / 2
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_BLANK])
    slide.shapes.add_picture(logo_fp, logo_left, logo_top, width = logo_width)
    return prs

def add_name_tag_slide(prs, company_name:str, product_tagline:str):
    """Function adding name and tag slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE])
    slide.placeholders[0].text = company_name
    slide.placeholders[1].text = product_tagline
    return prs

def make_presentation(logo_fp:str, company_name:str, product_tagline:str, save_fp:str):
    """Function making pitch deck"""
    prs = Presentation()
    add_logo_slide(prs, logo_fp)
    add_name_tag_slide(prs, company_name, product_tagline)
    prs.save(save_fp)

with open("results/Submarines-for-cats/202402081641/text.txt", mode = "r") as f: 
    data = f.read()
text = json.loads(data)
name = text.get("name")
tag = text.get("tag")
make_presentation("results/Submarines-for-cats/202402081641/logo.jpg", name, tag,
                  "results/Submarines-for-cats/202402081641/prs.pptx")


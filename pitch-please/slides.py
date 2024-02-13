"""Module with util functions for working with pptx Python library"""

from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from urllib.request import urlopen
import json

SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_CONTENT = 1
SLIDE_LAYOUT_SECTION_HEADER = 2
SLIDE_LAYOUT_TWO_CONTENT = 3
SLIDE_LAYOUT_COMPARISON = 4
SLIDE_LAYOUT_TITLE_ONLY = 5
SLIDE_LAYOUT_BLANK = 6
SLIDE_LAYOUT_CONTENT_CAPTION = 7
SLIDE_LAYOUT_PICTURE_CAPTION = 8

def add_logo_slide(prs, logo_url:str):
    """Function adding logo slide to presentation"""
    logo_data = BytesIO(urlopen(logo_url).read())
    logo_width = logo_height = Inches(6)
    logo_left = (prs.slide_width - logo_width) / 2
    logo_top = (prs.slide_height - logo_height) / 2
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_BLANK])
    slide.shapes.add_picture(logo_data, logo_left, logo_top, width = logo_width)
    return prs

def add_name_tag_slide(prs, company_name:str, product_tagline:str):
    """Function adding name and tag slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE])
    slide.placeholders[0].text = company_name
    slide.placeholders[1].text = product_tagline
    return prs

def add_problem_slide(prs, problems:list):
    """Function adding problem slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_CONTENT])
    slide.shapes.title.text = "The Problem Today"
    tf = slide.shapes.placeholders[1].text_frame
    for problem in problems:
        tf.add_paragraph().text = problem
    return prs

def add_solution_slide(prs, solution:str):
    """Function adding solution slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_CONTENT])
    slide.shapes.title.text = "The Solution"
    slide.shapes.placeholders[1].text_frame.paragraphs[0].text = solution
    return prs

def make_presentation(text:dict, logo_url:str):
    """Function making pitch deck"""
    prs = Presentation()
    add_name_tag_slide(prs, text["name"], text["tag"])
    add_problem_slide(prs, text["problem"])
    add_solution_slide(prs, text["solution"])
    add_logo_slide(prs, logo_url)
    return prs
